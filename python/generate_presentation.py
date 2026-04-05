"""
Generate Executive Presentation Deck
Highlights Churn Risk customers and overall segments
"""

import os
import sys
import pandas as pd
from sqlalchemy import create_engine
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Add python folder to path so imports work
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from config import CONNECTION_STRING, PROJECT_ROOT

OUTPUT_DIR = os.path.join(PROJECT_ROOT, "powerbi")
os.makedirs(OUTPUT_DIR, exist_ok=True)
PPT_PATH = os.path.join(OUTPUT_DIR, "Consumer360_Executive_Deck.pptx")

# Connect and get data
engine = create_engine(CONNECTION_STRING)
try:
    df = pd.read_sql("SELECT * FROM RFM_Results_Output", engine)
except Exception as e:
    print(f"Error reading from database: {e}")
    print("Trying to run without database info or you may need to run the pipeline first.")
    sys.exit(1)
finally:
    engine.dispose()

# Process data
total_customers = len(df)
total_revenue = df['TotalSpend'].sum() if 'TotalSpend' in df.columns else 0
avg_clv = df['CLV_Predicted'].mean() if 'CLV_Predicted' in df.columns else 0

# At Risk Customers
churn_risk = df[df['Segment'].isin(['At Risk', 'Hibernating'])].copy()
if churn_risk.empty:
    # Fallback to some logic if there are no "At Risk" literally named
    churn_risk = df[df['RFM_Score'] < 2.5]
    
churn_risk = churn_risk.sort_values(by='TotalSpend', ascending=False).head(10)

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Consumer360 RFM Analytics"
    subtitle.text = "Executive Summary & Churn Risk Analysis\nGenerated Automatically"
    
def add_summary_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = f"Total Customers Analyzed: {total_customers:,}"
    p.font.size = Pt(24)
    
    p2 = tf.add_paragraph()
    p2.text = f"Total Revenue Processed: ${total_revenue:,.2f}"
    p2.font.size = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = f"Average Predicted CLV (5-Year): ${avg_clv:,.2f}"
    p3.font.size = Pt(24)
    
    p4 = tf.add_paragraph()
    p4.text = f"High Risk Customers Identified: {len(churn_risk)} top customers at risk"
    p4.font.size = Pt(24)
    p4.font.color.rgb = RGBColor(255, 0, 0)
    
def add_churn_risk_slide(prs):
    slide_layout = prs.slide_layouts[5] # blank with title
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Actionable List: Top Churn Risk Customers"
    
    rows = len(churn_risk) + 1
    cols = 5
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Customer Name", "Segment", "RFM Score", "Total Spend", "Predicted CLV"]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                
    # Data
    for row_idx, (_, row) in enumerate(churn_risk.iterrows(), start=1):
        table.cell(row_idx, 0).text = str(row.get('CustomerName', 'Unknown'))
        table.cell(row_idx, 1).text = str(row.get('Segment', 'Unknown'))
        table.cell(row_idx, 2).text = f"{row.get('RFM_Score', 0):.2f}"
        table.cell(row_idx, 3).text = f"${row.get('TotalSpend', 0):,.2f}"
        table.cell(row_idx, 4).text = f"${row.get('CLV_Predicted', 0):,.2f}"
        
        for col_idx in range(cols):
            for paragraph in table.cell(row_idx, col_idx).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

# Create presentation
prs = Presentation()
add_title_slide(prs)
add_summary_slide(prs)
add_churn_risk_slide(prs)

prs.save(PPT_PATH)
print("="*60)
print(f"SUCCESS! Executive Presentation Deck generated at:\n{PPT_PATH}")
print("="*60)
